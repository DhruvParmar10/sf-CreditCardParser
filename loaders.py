"""
Universal document loader factory with comprehensive format support.
Gracefully handles multiple document formats with fallback strategies.
"""

import os
import tempfile
import zipfile
import email
from io import BytesIO
from typing import List, Optional, Union, Tuple
from pathlib import Path

# Document processing imports
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.document_loaders import PyMuPDFLoader
from docx import Document as DocxDocument
from pptx import Presentation
import pytesseract
from PIL import Image
import xml.etree.ElementTree as ET
import requests

# Enhanced imports
import magic
import logging

logger = logging.getLogger(__name__)

class DocumentLoader:
    """Universal document loader with comprehensive format support."""
    
    def __init__(self, max_recursion_depth: int = 3):
        self.max_recursion_depth = max_recursion_depth
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200,
            length_function=len,
        )
        
        # Supported MIME types mapping
        self.mime_handlers = {
            'application/pdf': self._load_pdf,
            'application/vnd.openxmlformats-officedocument.wordprocessingml.document': self._load_docx,
            'application/vnd.openxmlformats-officedocument.presentationml.presentation': self._load_pptx,
            'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': self._load_xlsx,
            'application/zip': self._load_zip,
            'application/x-zip-compressed': self._load_zip,
            'message/rfc822': self._load_email,
            'text/xml': self._load_xml,
            'application/xml': self._load_xml,
            'image/jpeg': self._load_image,
            'image/png': self._load_image,
            'image/tiff': self._load_image,
            'text/plain': self._load_text,
            'text/html': self._load_html,
        }
        
        # File extension fallbacks
        self.extension_handlers = {
            '.pdf': self._load_pdf,
            '.docx': self._load_docx,
            '.pptx': self._load_pptx,
            '.xlsx': self._load_xlsx,
            '.zip': self._load_zip,
            '.eml': self._load_email,
            '.xml': self._load_xml,
            '.jpg': self._load_image,
            '.jpeg': self._load_image,
            '.png': self._load_image,
            '.tiff': self._load_image,
            '.tif': self._load_image,
            '.txt': self._load_text,
            '.html': self._load_html,
            '.htm': self._load_html,
        }

    def load_file(self, file_path: Union[str, Path], url: Optional[str] = None, recursion_depth: int = 0) -> Tuple[str, str]:
        """
        Universal file loader with automatic format detection.
        
        Args:
            file_path: Path to the file or temp file path for uploaded content
            url: Optional URL for context (if file was downloaded)
            recursion_depth: Current recursion depth to prevent infinite loops
            
        Returns:
            Tuple of (extracted_text, detected_format)
        """
        # Prevent excessive recursion (especially for nested ZIP files)
        if recursion_depth >= self.max_recursion_depth:
            logger.warning(f"Maximum recursion depth ({self.max_recursion_depth}) reached for {file_path}")
            return f"Skipped due to recursion limit: {Path(file_path).name}", "recursion_limit"
        
        try:
            file_path = Path(file_path)
            
            # Detect MIME type
            mime_type = self._detect_mime_type(file_path)
            logger.info(f"Detected MIME type: {mime_type} for file: {file_path.name} (extension: {file_path.suffix}) [depth: {recursion_depth}]")
            
            # Try MIME type handler first
            if mime_type in self.mime_handlers:
                logger.info(f"Using MIME handler for {mime_type}")
                if mime_type in ['application/zip', 'application/x-zip-compressed']:
                    # Pass recursion depth to ZIP handler
                    text = self._load_zip(file_path, recursion_depth)
                else:
                    text = self.mime_handlers[mime_type](file_path)
                return text, mime_type
            
            # Fallback to extension-based detection
            extension = file_path.suffix.lower()
            logger.info(f"Trying extension-based detection for: {extension}")
            if extension in self.extension_handlers:
                logger.info(f"Using extension handler for {extension}")
                if extension == '.zip':
                    # Pass recursion depth to ZIP handler
                    text = self._load_zip(file_path, recursion_depth)
                else:
                    text = self.extension_handlers[extension](file_path)
                return text, f"extension:{extension}"
            
            # Final fallback - try as text
            logger.warning(f"Unknown format for {file_path.name} (MIME: {mime_type}, ext: {extension}), attempting text extraction")
            text = self._load_text(file_path)
            return text, "text/plain"
            
        except Exception as e:
            logger.error(f"Error loading file {file_path}: {str(e)}")
            return f"Error processing file: {str(e)}", "error"

    def _detect_mime_type(self, file_path: Path) -> str:
        """Detect MIME type using python-magic with enhanced fallbacks."""
        # First try: python-magic
        try:
            mime_type = magic.from_file(str(file_path), mime=True)
            # Accept the result if it's not the generic fallback
            if mime_type and mime_type not in ['application/octet-stream', 'text/plain']:
                return mime_type
        except Exception as e:
            logger.warning(f"MIME detection failed for {file_path}: {e}")
        
        # Enhanced fallback 1: Check file signature/header
        try:
            with open(file_path, 'rb') as f:
                header = f.read(32)  # Read more bytes for better detection
                
            # Common file signatures with more specific detection
            if header.startswith(b'%PDF'):
                return 'application/pdf'
            elif header.startswith(b'PK\x03\x04'):
                # ZIP-based formats - check further into the file
                try:
                    with zipfile.ZipFile(file_path, 'r') as zf:
                        file_list = zf.namelist()
                        # Check for Office document patterns
                        if any('word/' in f for f in file_list):
                            return 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'
                        elif any('ppt/' in f for f in file_list):
                            return 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
                        elif any('xl/' in f for f in file_list):
                            return 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
                        else:
                            # Check file extension as secondary indicator
                            ext = file_path.suffix.lower()
                            if ext == '.docx':
                                return 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'
                            elif ext == '.pptx':
                                return 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
                            else:
                                return 'application/zip'
                except:
                    return 'application/zip'
            elif header.startswith(b'\xff\xd8\xff'):
                return 'image/jpeg'
            elif header.startswith(b'\x89PNG'):
                return 'image/png'
            elif header.startswith(b'GIF8'):
                return 'image/gif'
            elif header.startswith(b'II*\x00') or header.startswith(b'MM\x00*'):
                return 'image/tiff'
            elif header.startswith(b'<?xml'):
                return 'text/xml'
            elif header.startswith(b'<!DOCTYPE html') or header.startswith(b'<html'):
                return 'text/html'
        except Exception:
            pass
        
        # Fallback 2: Extension-based detection
        extension_mime_map = {
            '.pdf': 'application/pdf',
            '.docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
            '.pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
            '.xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
            '.zip': 'application/zip',
            '.eml': 'message/rfc822',
            '.xml': 'text/xml',
            '.jpg': 'image/jpeg',
            '.jpeg': 'image/jpeg',
            '.png': 'image/png',
            '.gif': 'image/gif',
            '.tiff': 'image/tiff',
            '.tif': 'image/tiff',
            '.txt': 'text/plain',
            '.html': 'text/html',
            '.htm': 'text/html',
        }
        return extension_mime_map.get(file_path.suffix.lower(), 'application/octet-stream')

    def _load_pdf(self, file_path: Path) -> str:
        """Load PDF documents using PyMuPDF."""
        loader = PyMuPDFLoader(str(file_path))
        docs = loader.load()
        return "\n\n".join([doc.page_content for doc in docs])

    def _load_docx(self, file_path: Path) -> str:
        """Load Word documents."""
        doc = DocxDocument(str(file_path))
        text_content = []
        
        # Extract paragraphs
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text_content.append(paragraph.text.strip())
        
        # Extract tables
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join([cell.text.strip() for cell in row.cells])
                if row_text.strip():
                    text_content.append(row_text)
        
        return "\n\n".join(text_content)

    def _load_pptx(self, file_path: Path) -> str:
        """Load PowerPoint presentations."""
        prs = Presentation(str(file_path))
        text_content = []
        
        for i, slide in enumerate(prs.slides, 1):
            slide_content = [f"=== Slide {i} ==="]
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_content.append(shape.text.strip())
            
            text_content.append("\n".join(slide_content))
        
        return "\n\n".join(text_content)

    def _load_xlsx(self, file_path: Path) -> str:
        """Load Excel spreadsheets."""
        try:
            import openpyxl
            workbook = openpyxl.load_workbook(str(file_path), data_only=True)
            text_content = []
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                sheet_content = [f"=== Sheet: {sheet_name} ==="]
                
                # Extract data from cells
                for row in sheet.iter_rows(values_only=True):
                    # Filter out empty rows and format as tab-separated
                    if any(cell is not None for cell in row):
                        row_text = "\t".join(str(cell) if cell is not None else "" for cell in row)
                        if row_text.strip():
                            sheet_content.append(row_text)
                
                if len(sheet_content) > 1:  # Has content beyond header
                    text_content.append("\n".join(sheet_content))
            
            workbook.close()
            return "\n\n".join(text_content)
            
        except ImportError:
            logger.error("openpyxl not installed - cannot process Excel files")
            return "Error: openpyxl library required for Excel file processing"
        except Exception as e:
            logger.error(f"Excel processing error: {e}")
            return f"Error processing Excel file: {str(e)}"

    def _load_zip(self, file_path: Path, recursion_depth: int = 0) -> str:
        """Load and process ZIP archives with recursion protection."""
        text_content = []
        processed_files = 0
        max_files_per_zip = 20  # Limit files processed per ZIP to prevent resource exhaustion
        
        # Check for potential ZIP bomb (many files of similar size)
        with zipfile.ZipFile(file_path, 'r') as zip_ref:
            file_list = zip_ref.filelist
            zip_files = [f for f in file_list if f.filename.lower().endswith('.zip')]
            
            # If this ZIP contains mostly other ZIP files, it might be a ZIP bomb
            if len(zip_files) > 10 and len(zip_files) / len(file_list) > 0.8:
                logger.warning(f"Potential ZIP bomb detected in {file_path.name}: {len(zip_files)} ZIP files out of {len(file_list)} total files")
                if recursion_depth > 0:
                    return f"ZIP bomb detected - contains {len(zip_files)} nested ZIP files. Processing stopped for security."
                else:
                    # For the top level, try to process just a few files
                    max_files_per_zip = 3
        
        with zipfile.ZipFile(file_path, 'r') as zip_ref:
            for file_info in zip_ref.filelist:
                if file_info.is_dir():
                    continue
                
                # Limit number of files processed per ZIP
                if processed_files >= max_files_per_zip:
                    remaining = len(zip_ref.filelist) - processed_files
                    logger.warning(f"Reached max files limit ({max_files_per_zip}) for ZIP: {file_path.name}")
                    text_content.append(f"... and {remaining} more files (truncated for security)")
                    break
                
                try:
                    # Get the original file extension
                    original_path = Path(file_info.filename)
                    file_extension = original_path.suffix
                    
                    # Skip nested ZIP files if we're already in recursion or if it looks like a ZIP bomb
                    if file_extension.lower() == '.zip':
                        if recursion_depth >= 1:
                            logger.info(f"Skipping nested ZIP file {file_info.filename} at depth {recursion_depth}")
                            text_content.append(f"=== {file_info.filename} (skipped - nested ZIP) ===\nNested ZIP file skipped to prevent recursion")
                            processed_files += 1
                            continue
                        else:
                            # Even at top level, limit ZIP processing if it seems suspicious
                            if len([f for f in zip_ref.filelist if f.filename.lower().endswith('.zip')]) > 5:
                                logger.warning(f"Many ZIP files detected, limiting processing of {file_info.filename}")
                                text_content.append(f"=== {file_info.filename} (limited processing) ===\nNested ZIP files detected - limited processing for security")
                                processed_files += 1
                                continue
                    
                    # Skip if no extension and apply common defaults
                    if not file_extension:
                        # Try to detect from file content or apply reasonable defaults
                        file_data = zip_ref.read(file_info.filename)
                        if file_data.startswith(b'%PDF'):
                            file_extension = '.pdf'
                        elif file_data.startswith(b'PK\x03\x04'):
                            # Could be docx, pptx, xlsx, or other zip-based format
                            if 'word' in file_info.filename.lower() or 'document' in file_info.filename.lower():
                                file_extension = '.docx'
                            elif 'ppt' in file_info.filename.lower() or 'presentation' in file_info.filename.lower():
                                file_extension = '.pptx'
                            elif 'xl' in file_info.filename.lower() or 'sheet' in file_info.filename.lower() or 'excel' in file_info.filename.lower():
                                file_extension = '.xlsx'
                            else:
                                file_extension = '.zip'
                        elif file_data.startswith((b'\xff\xd8\xff', b'\x89PNG')):
                            file_extension = '.jpg' if file_data.startswith(b'\xff\xd8\xff') else '.png'
                        else:
                            file_extension = '.txt'  # Default fallback
                        
                        logger.info(f"Auto-detected extension {file_extension} for {file_info.filename}")
                    
                    # Extract to temporary file with proper extension
                    with tempfile.NamedTemporaryFile(delete=False, suffix=file_extension) as temp_file:
                        temp_file.write(zip_ref.read(file_info.filename))
                        temp_path = Path(temp_file.name)
                    
                    logger.info(f"Processing {file_info.filename} as {temp_path} with extension {file_extension} [depth: {recursion_depth}]")
                    
                    # Recursively process extracted file with incremented depth
                    content, detected_format = self.load_file(temp_path, recursion_depth=recursion_depth + 1)
                    text_content.append(f"=== {file_info.filename} ({detected_format}) ===\n{content}")
                    
                    # Clean up
                    os.unlink(temp_path)
                    processed_files += 1
                    
                except Exception as e:
                    logger.warning(f"Error processing {file_info.filename} from ZIP: {e}")
                    processed_files += 1
                    continue
        
        if not text_content:
            return "ZIP file processed but no extractable content found (may be a ZIP bomb or corrupted archive)"
        
        return "\n\n".join(text_content)

    def _load_email(self, file_path: Path) -> str:
        """Load email files (.eml)."""
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            msg = email.message_from_file(f)
        
        text_content = []
        
        # Extract headers
        headers = [
            f"From: {msg.get('From', 'Unknown')}",
            f"To: {msg.get('To', 'Unknown')}",
            f"Subject: {msg.get('Subject', 'No Subject')}",
            f"Date: {msg.get('Date', 'Unknown')}",
        ]
        text_content.append("=== Email Headers ===\n" + "\n".join(headers))
        
        # Extract body
        if msg.is_multipart():
            for part in msg.walk():
                if part.get_content_type() == "text/plain":
                    body = part.get_payload(decode=True)
                    if body:
                        text_content.append("=== Email Body ===\n" + body.decode('utf-8', errors='ignore'))
        else:
            body = msg.get_payload(decode=True)
            if body:
                text_content.append("=== Email Body ===\n" + body.decode('utf-8', errors='ignore'))
        
        return "\n\n".join(text_content)

    def _load_xml(self, file_path: Path) -> str:
        """Load XML documents."""
        try:
            tree = ET.parse(file_path)
            root = tree.getroot()
            
            # Extract all text content
            text_content = [f"=== XML Document: {root.tag} ==="]
            
            def extract_text(element, level=0):
                indent = "  " * level
                if element.text and element.text.strip():
                    text_content.append(f"{indent}{element.tag}: {element.text.strip()}")
                
                for child in element:
                    extract_text(child, level + 1)
            
            extract_text(root)
            return "\n".join(text_content)
            
        except ET.ParseError as e:
            logger.error(f"XML parsing error: {e}")
            return self._load_text(file_path)

    def _load_image(self, file_path: Path) -> str:
        """Load images using OCR with enhanced preprocessing."""
        try:
            image = Image.open(file_path)
            
            # Enhanced OCR with multiple configurations
            configs = [
                '--psm 1 --oem 3',  # Automatic page segmentation with orientation
                '--psm 3 --oem 3',  # Fully automatic page segmentation (default)
                '--psm 6 --oem 3',  # Uniform block of text
                '--psm 4 --oem 3',  # Single column of text
                '--psm 8 --oem 3',  # Single word
                '--psm 13 --oem 3', # Raw line (treat as single text line)
            ]
            
            best_text = ""
            max_confidence = 0
            
            for config in configs:
                try:
                    # Try with original image
                    text = pytesseract.image_to_string(image, config=config)
                    confidence = len(text.strip())  # Simple confidence metric
                    
                    if confidence > max_confidence:
                        max_confidence = confidence
                        best_text = text
                        
                except Exception as e:
                    logger.warning(f"OCR config {config} failed: {e}")
                    continue
            
            # If no good result, try with image preprocessing
            if len(best_text.strip()) < 50:
                try:
                    # Convert to grayscale and enhance contrast
                    gray_image = image.convert('L')
                    enhanced_image = gray_image.point(lambda p: p * 1.2)
                    best_text = pytesseract.image_to_string(enhanced_image, config='--psm 3 --oem 3')
                except Exception as e:
                    logger.warning(f"Enhanced OCR failed: {e}")
            
            return best_text.strip() if best_text.strip() else "No text could be extracted from image"
            
        except Exception as e:
            logger.error(f"Image processing error: {e}")
            return f"Error processing image: {str(e)}"

    def _load_text(self, file_path: Path) -> str:
        """Load plain text files."""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                return f.read()
        except Exception as e:
            logger.error(f"Text loading error: {e}")
            return f"Error reading text file: {str(e)}"

    def _load_html(self, file_path: Path) -> str:
        """Load HTML files and extract text content."""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                html_content = f.read()
            
            # Simple HTML text extraction (could be enhanced with BeautifulSoup)
            import re
            # Remove script and style elements
            html_content = re.sub(r'<script.*?</script>', '', html_content, flags=re.DOTALL | re.IGNORECASE)
            html_content = re.sub(r'<style.*?</style>', '', html_content, flags=re.DOTALL | re.IGNORECASE)
            # Remove HTML tags
            text_content = re.sub(r'<[^>]+>', '', html_content)
            # Clean up whitespace
            text_content = re.sub(r'\s+', ' ', text_content).strip()
            
            return text_content
            
        except Exception as e:
            logger.error(f"HTML processing error: {e}")
            return self._load_text(file_path)

    def download_and_load(self, url: str) -> Tuple[str, str]:
        """Download file from URL and load it."""
        try:
            response = requests.get(url, timeout=30)
            response.raise_for_status()
            
            # Determine file extension from URL or content type
            file_extension = ""
            
            # First try to get extension from URL
            from urllib.parse import urlparse, unquote
            parsed_url = urlparse(unquote(url))
            url_path = Path(parsed_url.path)
            if url_path.suffix:
                file_extension = url_path.suffix
            
            # If no extension from URL, try content type
            if not file_extension:
                content_type = response.headers.get('content-type', '').split(';')[0].strip()
                content_type_extensions = {
                    'application/pdf': '.pdf',
                    'application/vnd.openxmlformats-officedocument.wordprocessingml.document': '.docx',
                    'application/vnd.openxmlformats-officedocument.presentationml.presentation': '.pptx',
                    'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': '.xlsx',
                    'application/zip': '.zip',
                    'image/jpeg': '.jpg',
                    'image/png': '.png',
                    'text/plain': '.txt',
                    'text/html': '.html',
                }
                file_extension = content_type_extensions.get(content_type, '')
            
            # Create temporary file with proper extension
            with tempfile.NamedTemporaryFile(delete=False, suffix=file_extension) as temp_file:
                temp_file.write(response.content)
                temp_path = Path(temp_file.name)
            
            logger.info(f"Downloaded {len(response.content)} bytes to {temp_path} with extension {file_extension}")
            
            # Load the downloaded file
            text, format_type = self.load_file(temp_path, url)
            
            # Clean up
            os.unlink(temp_path)
            
            return text, format_type
            
        except Exception as e:
            logger.error(f"Error downloading from URL {url}: {e}")
            return f"Error downloading file: {str(e)}", "error"

    def split_text(self, text: str) -> List[str]:
        """Split text into chunks for processing."""
        return self.text_splitter.split_text(text)


# Factory function for easy usage
def load_file(file_path: Union[str, Path], url: Optional[str] = None) -> Tuple[str, str]:
    """
    Universal file loader factory function.
    
    Args:
        file_path: Path to the file to load
        url: Optional URL if file was downloaded
        
    Returns:
        Tuple of (extracted_text, detected_format)
    """
    loader = DocumentLoader()
    return loader.load_file(file_path, url, recursion_depth=0)


def download_and_load(url: str) -> Tuple[str, str]:
    """
    Download and load file from URL.
    
    Args:
        url: URL to download file from
        
    Returns:
        Tuple of (extracted_text, detected_format)
    """
    loader = DocumentLoader()
    return loader.download_and_load(url)
